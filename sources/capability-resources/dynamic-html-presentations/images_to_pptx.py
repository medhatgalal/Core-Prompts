#!/usr/bin/env python3
"""Build a visually faithful image-only 16:9 PPTX from slide-*.png files."""

from __future__ import annotations

import argparse
import re
import sys
import tempfile
from pathlib import Path

try:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as exc:
    print(
        "error: python-pptx==1.0.2 is required. Install it temporarily with:\n"
        "  DEPS=\"$(mktemp -d)\"\n"
        "  python3 -m pip install --target \"$DEPS\" 'python-pptx==1.0.2'\n"
        "  PYTHONPATH=\"$DEPS\" python3 images_to_pptx.py SLIDES_DIR OUTPUT.pptx",
        file=sys.stderr,
    )
    raise SystemExit(2) from exc

SLIDE_PATTERN = re.compile(r"^slide-(\d+)\.png$", re.IGNORECASE)
WIDESCREEN_WIDTH_INCHES = 13.333333
WIDESCREEN_HEIGHT_INCHES = 7.5


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Create a widescreen, image-only PPTX from naturally sorted slide-*.png files."
    )
    parser.add_argument("image_dir", type=Path, help="Directory containing slide-*.png files")
    parser.add_argument("output", type=Path, help="Destination .pptx file")
    parser.add_argument("--title", default="", help="Optional presentation title metadata")
    parser.add_argument(
        "--force",
        action="store_true",
        help="Atomically replace an existing output PPTX after the new file verifies",
    )
    return parser.parse_args()


def discover_images(image_dir: Path) -> list[Path]:
    if not image_dir.is_dir():
        raise ValueError(f"Image directory does not exist: {image_dir}")

    numbered: list[tuple[int, Path]] = []
    for path in image_dir.iterdir():
        match = SLIDE_PATTERN.match(path.name)
        if match and path.is_file():
            numbered.append((int(match.group(1)), path))

    if not numbered:
        raise ValueError(f"No slide-*.png files found in {image_dir}")

    numbered.sort(key=lambda item: (item[0], item[1].name.lower()))
    numbers = [number for number, _ in numbered]
    expected = list(range(1, len(numbers) + 1))
    if numbers != expected:
        raise ValueError(
            f"Slide numbering must be contiguous from 1; found {numbers}, expected {expected}"
        )
    return [path for _, path in numbered]


def validate_dimensions(images: list[Path]) -> tuple[int, int]:
    dimensions: list[tuple[Path, tuple[int, int]]] = []
    for path in images:
        if path.stat().st_size <= 0:
            raise ValueError(f"Image is empty: {path}")
        with Image.open(path) as image:
            image.verify()
        with Image.open(path) as image:
            dimensions.append((path, image.size))

    width, height = dimensions[0][1]
    mismatched = [(path, size) for path, size in dimensions if size != (width, height)]
    if mismatched:
        details = ", ".join(f"{path.name}={size[0]}x{size[1]}" for path, size in mismatched)
        raise ValueError(
            f"All images must use the same dimensions; expected {width}x{height}, found {details}"
        )
    if width * 9 != height * 16:
        raise ValueError(f"Images must be exactly 16:9; found {width}x{height}")
    return width, height


def build_presentation(images: list[Path], output: Path, title: str, force: bool) -> None:
    output = output.expanduser().resolve()
    if output.suffix.lower() != ".pptx":
        raise ValueError(f"Output must end in .pptx: {output}")
    output.parent.mkdir(parents=True, exist_ok=True)
    if output.exists() and not force:
        raise ValueError(f"Output already exists; pass --force to replace it: {output}")

    presentation = Presentation()
    presentation.slide_width = Inches(WIDESCREEN_WIDTH_INCHES)
    presentation.slide_height = Inches(WIDESCREEN_HEIGHT_INCHES)
    if title:
        presentation.core_properties.title = title

    blank_layout = presentation.slide_layouts[6]
    for image_path in images:
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path.resolve()),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    temporary: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            prefix=f".{output.stem}-",
            suffix=".pptx",
            dir=output.parent,
            delete=False,
        ) as handle:
            temporary = Path(handle.name)

        presentation.save(temporary)
        if temporary.stat().st_size <= 0:
            raise RuntimeError(f"Temporary PPTX was not created correctly: {temporary}")

        verification = Presentation(temporary)
        if len(verification.slides) != len(images):
            raise RuntimeError(
                f"PPTX verification failed: {len(verification.slides)} slides for {len(images)} images"
            )
        if output.exists() and not force:
            raise ValueError(f"Output appeared during generation; pass --force to replace it: {output}")
        temporary.replace(output)
        temporary = None
    finally:
        if temporary is not None:
            temporary.unlink(missing_ok=True)


def main() -> int:
    args = parse_args()
    try:
        image_dir = args.image_dir.expanduser().resolve()
        images = discover_images(image_dir)
        width, height = validate_dimensions(images)
        build_presentation(images, args.output, args.title, args.force)
    except (OSError, ValueError, RuntimeError) as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 1

    output = args.output.expanduser().resolve()
    print(
        f"Created and verified {output} with {len(images)} image-only slide(s) "
        f"from {width}x{height} PNGs."
    )
    print("Note: slide content is visually faithful but flattened and not independently editable.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
